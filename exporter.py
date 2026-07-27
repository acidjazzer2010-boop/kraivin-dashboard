import os
import io
import smtplib
from email.mime.multipart import MIMEMultipart
from email.mime.base import MIMEBase
from email.mime.text import MIMEText
from email import encoders

import numpy as np
import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_chart_liquidity(x_labels, cash_balance):
    fig, ax = plt.subplots(figsize=(10, 4.2), dpi=200)
    ax.plot(x_labels, cash_balance, color='#642A38', linewidth=3, marker='o', markersize=5)
    ax.fill_between(x_labels, cash_balance, 0, color='#642A38', alpha=0.1)
    ax.axhline(0, color='red', linestyle='--', linewidth=1.5, label='Критический уровень (0)')
    ax.set_title("Динамика ликвидности и остаток средств", fontsize=13, fontweight='bold', color='#642A38', pad=12)
    ax.set_ylabel("Рубли", fontsize=10, color='#333333')
    plt.xticks(rotation=25, ha='right', fontsize=9)
    plt.grid(axis='y', linestyle=':', alpha=0.5)
    ax.legend(frameon=True, facecolor='white', edgecolor='none')
    plt.tight_layout()
    
    img_io = io.BytesIO()
    plt.savefig(img_io, format='png', bbox_inches='tight')
    plt.close(fig)
    img_io.seek(0)
    return img_io

def create_chart_cash_flow(x_labels, inflows, outflows, net_cf):
    fig, ax = plt.subplots(figsize=(10, 4.2), dpi=200)
    x_indices = np.arange(len(x_labels))
    width = 0.35
    
    ax.bar(x_indices - width/2, inflows, width, label='Поступления', color='#E3C293')
    ax.bar(x_indices + width/2, -outflows, width, label='Выплаты', color='#642A38')
    ax.plot(x_indices, net_cf, color='#B88645', linewidth=2.5, marker='o', label='Чистый поток')
    
    ax.set_xticks(x_indices)
    ax.set_xticklabels(x_labels, rotation=25, ha='right', fontsize=9)
    ax.set_title("Структура месячного денежного потока", fontsize=13, fontweight='bold', color='#642A38', pad=12)
    ax.set_ylabel("Рубли", fontsize=10, color='#333333')
    ax.legend(frameon=True, facecolor='white', edgecolor='none')
    plt.grid(axis='y', linestyle=':', alpha=0.5)
    plt.tight_layout()
    
    img_io = io.BytesIO()
    plt.savefig(img_io, format='png', bbox_inches='tight')
    plt.close(fig)
    img_io.seek(0)
    return img_io

def create_chart_margin_ebitda(x_labels, rev, opex, taxes_and_commissions, cogs_payments, margin_pct):
    fig, ax1 = plt.subplots(figsize=(10, 4.2), dpi=200)
    
    ebitda = rev - opex - taxes_and_commissions - (cogs_payments * 0.3)
    
    ax1.bar(x_labels, ebitda, color='#642A38', alpha=0.8, width=0.5, label='Операционная прибыль (EBITDA)')
    ax1.set_ylabel("EBITDA (руб)", color='#642A38', fontsize=10, fontweight='bold')
    ax1.tick_params(axis='y', labelcolor='#642A38')
    plt.xticks(rotation=25, ha='right', fontsize=9)
    
    ax2 = ax1.twinx()
    margin_dynamics = [margin_pct] * len(x_labels)
    ax2.plot(x_labels, margin_dynamics, color='#E3C293', linewidth=2.5, marker='s', label='Маржинальность (%)')
    ax2.set_ylabel("Маржа (%)", color='#B88645', fontsize=10, fontweight='bold')
    ax2.tick_params(axis='y', labelcolor='#B88645')
    ax2.set_ylim(0, 50)
    
    plt.title("Динамика маржинальности и операционной прибыли (EBITDA)", fontsize=13, fontweight='bold', color='#642A38', pad=12)
    plt.grid(axis='y', linestyle=':', alpha=0.3)
    plt.tight_layout()
    
    img_io = io.BytesIO()
    plt.savefig(img_io, format='png', bbox_inches='tight')
    plt.close(fig)
    img_io.seek(0)
    return img_io

def create_chart_sources(x_labels, rev, factoring_share):
    fig, ax = plt.subplots(figsize=(10, 4.2), dpi=200)
    x_indices = np.arange(len(x_labels))
    width = 0.5
    
    factoring_inflows = rev * 1.2 * (factoring_share / 100)
    direct_inflows = rev * 1.2 * ((100 - factoring_share) / 100)
    
    ax.bar(x_indices, direct_inflows, width, label='Оплата от клиентов', color='#642A38')
    ax.bar(x_indices, factoring_inflows, width, bottom=direct_inflows, label='Факторинг', color='#E3C293')
    
    ax.set_xticks(x_indices)
    ax.set_xticklabels(x_labels, rotation=25, ha='right', fontsize=9)
    ax.set_title("Структура притока денежных средств по источникам", fontsize=13, fontweight='bold', color='#642A38', pad=12)
    ax.set_ylabel("Рубли", fontsize=10, color='#333333')
    ax.legend(frameon=True, facecolor='white', edgecolor='none')
    plt.grid(axis='y', linestyle=':', alpha=0.5)
    plt.tight_layout()
    
    img_io = io.BytesIO()
    plt.savefig(img_io, format='png', bbox_inches='tight')
    plt.close(fig)
    img_io.seek(0)
    return img_io

def create_chart_cumulative(x_labels, cum_cf, initial_cash_buffer):
    fig, ax = plt.subplots(figsize=(10, 4.2), dpi=200)
    total_cash_line = cum_cf + initial_cash_buffer
    
    ax.plot(x_labels, total_cash_line, color='#642A38', linewidth=3, marker='o', label='Накопленный ДС с буфером')
    ax.axhline(initial_cash_buffer, color='#B88645', linestyle='--', linewidth=1.5, label='Стартовый буфер')
    ax.axhline(0, color='red', linestyle=':', linewidth=1.5, label='Нулевой баланс')
    ax.fill_between(x_labels, total_cash_line, initial_cash_buffer, where=(total_cash_line >= initial_cash_buffer), color='#642A38', alpha=0.1, interpolate=True)
    
    ax.set_title("Накопленный денежный поток (Cumulative Cash Flow) и зона безопасного остатка", fontsize=13, fontweight='bold', color='#642A38', pad=12)
    ax.set_ylabel("Рубли", fontsize=10, color='#333333')
    plt.xticks(rotation=25, ha='right', fontsize=9)
    ax.legend(frameon=True, facecolor='white', edgecolor='none')
    plt.grid(axis='y', linestyle=':', alpha=0.5)
    plt.tight_layout()
    
    img_io = io.BytesIO()
    plt.savefig(img_io, format='png', bbox_inches='tight')
    plt.close(fig)
    img_io.seek(0)
    return img_io

def create_chart_expenses_pie(sum_purchases, total_fot, total_taxes, total_other_opex):
    fig, ax = plt.subplots(figsize=(8, 4.2), dpi=200)
    labels = ['Закупки товара', 'ФОТ (Команда)', 'Налоги и сборы', 'Операционные расходы']
    sizes = [sum_purchases, total_fot, total_taxes, total_other_opex]
    colors = ['#642A38', '#E3C293', '#B88645', '#D0C2B8']
    
    ax.pie(sizes, labels=labels, colors=colors, autopct='%1.1f%%', startangle=140, textprops={'fontsize': 9})
    ax.set_title("Структура совокупных расходов", fontsize=13, fontweight='bold', color='#642A38', pad=12)
    plt.tight_layout()
    
    img_io = io.BytesIO()
    plt.savefig(img_io, format='png', bbox_inches='tight')
    plt.close(fig)
    img_io.seek(0)
    return img_io

def generate_pptx_bytes(period, start_month_idx, start_year, ru_months_full, x_labels, 
                        sum_rev, max_deficit, net_profit, roi, initial_cash_buffer, 
                        initial_purchase, inflows, outflows, cash_balance, rev, net_cf,
                        customer_delay_days, delay_days, factoring_share, logo_path,
                        opex=None, taxes_and_commissions=None, cogs_payments=None, cum_cf=None, 
                        sum_purchases=0, total_fot=0, total_taxes=0, total_other_opex=0, margin_pct=20):

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    c_wine = RGBColor(100, 42, 56)
    c_sand = RGBColor(227, 194, 147)
    c_dark = RGBColor(30, 30, 30)
    c_light_bg = RGBColor(248, 246, 244)
    c_white = RGBColor(255, 255, 255)
    c_card_border = RGBColor(220, 210, 205)

    font_black = "ua-BRAND-black"
    font_bold = "ua-BRAND-bold"
    font_regular = "ua-BRAND-regular"

    blank_layout = prs.slide_layouts[6]

    def add_corner_logo(slide):
        if logo_path and os.path.exists(logo_path):
            slide.shapes.add_picture(logo_path, Inches(11.8), Inches(0.4), width=Inches(0.9))

    def format_rub(val):
        return f"{val:,.0f}".replace(",", " ") + " руб."

    # СЛАЙД 1: ТИТУЛЬНЫЙ
    slide1 = prs.slides.add_slide(blank_layout)
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = c_wine
    bg1.line.fill.background()

    accent1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.0), Inches(0.15), Inches(3.5))
    accent1.fill.solid()
    accent1.fill.fore_color.rgb = c_sand
    accent1.line.fill.background()

    if logo_path and os.path.exists(logo_path):
        slide1.shapes.add_picture(logo_path, Inches(1.2), Inches(0.8), width=Inches(1.8))

    txBox = slide1.shapes.add_textbox(Inches(1.2), Inches(2.5), Inches(10.5), Inches(3.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "КРАЙВИН"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = c_white
    p.font.name = font_black

    p2 = tf.add_paragraph()
    p2.text = "Финансовая модель и углубленный анализ денежных потоков"
    p2.font.size = Pt(26)
    p2.font.color.rgb = c_sand
    p2.font.name = font_bold
    p2.space_before = Pt(10)

    p3 = tf.add_paragraph()
    p3.text = f"Горизонт планирования: {period} мес.  |  Старт: {ru_months_full[start_month_idx]} {start_year}"
    p3.font.size = Pt(16)
    p3.font.color.rgb = c_white
    p3.font.name = font_regular
    p3.space_before = Pt(35)

    # СЛАЙД 2: KPI КАРТОЧКИ
    slide2 = prs.slides.add_slide(blank_layout)
    bg2 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg2.fill.solid()
    bg2.fill.fore_color.rgb = c_light_bg
    bg2.line.fill.background()
    add_corner_logo(slide2)

    title_box2 = slide2.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(10.0), Inches(0.8))
    p_t2 = title_box2.text_frame.paragraphs[0]
    p_t2.text = "Ключевые финансовые результаты"
    p_t2.font.size = Pt(26)
    p_t2.font.bold = True
    p_t2.font.color.rgb = c_wine
    p_t2.font.name = font_bold

    kpis = [
        ("Суммарная выручка", format_rub(sum_rev)),
        ("Макс. кассовый разрыв", format_rub(max_deficit)),
        ("Чистая прибыль", format_rub(net_profit)),
        ("Рентабельность по ЧП", f"{roi:.1f}%"),
        ("Начальный буфер ДС", format_rub(initial_cash_buffer)),
        ("Первоначальная закупка", format_rub(initial_purchase))
    ]

    card_w, card_h = Inches(3.6), Inches(2.1)
    start_x, start_y = Inches(0.8), Inches(1.5)
    gap_x, gap_y = Inches(0.4), Inches(0.35)

    for idx, (label, val) in enumerate(kpis):
        col_idx = idx % 3
        row_idx = idx // 3
        x = start_x + col_idx * (card_w + gap_x)
        y = start_y + row_idx * (card_h + gap_y)

        card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = c_white
        card.line.color.rgb = c_card_border
        card.line.width = Pt(1)

        tf_card = card.text_frame
        tf_card.word_wrap = True
        tf_card.margin_left = Inches(0.25)
        tf_card.margin_top = Inches(0.2)
        
        p_lbl = tf_card.paragraphs[0]
        p_lbl.text = label.upper()
        p_lbl.font.size = Pt(11)
        p_lbl.font.color.rgb = RGBColor(120, 120, 120)
        p_lbl.font.bold = True
        p_lbl.font.name = font_bold

        p_val = tf_card.add_paragraph()
        p_val.text = val
        p_val.font.size = Pt(20)
        p_val.font.color.rgb = c_wine
        p_val.font.bold = True
        p_val.font.name = font_black
        p_val.space_before = Pt(8)

    # Добавление слайдов с графиками
    slides_data = [
        ("Динамика ликвидности и остаток средств", create_chart_liquidity(x_labels, cash_balance)),
        ("Структура месячного денежного потока", create_chart_cash_flow(x_labels, inflows, outflows, net_cf)),
        ("Динамика маржинальности и операционной прибыли (EBITDA)", create_chart_margin_ebitda(x_labels, rev, opex, taxes_and_commissions, cogs_payments, margin_pct)),
        ("Структура притока денежных средств по источникам", create_chart_sources(x_labels, rev, factoring_share)),
        ("Накопленный денежный поток и зона безопасного остатка", create_chart_cumulative(x_labels, cum_cf, initial_cash_buffer)),
        ("Структура совокупных расходов", create_chart_expenses_pie(sum_purchases, total_fot, total_taxes, total_other_opex))
    ]

    for title, img_func in slides_data:
        slide = prs.slides.add_slide(blank_layout)
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = c_light_bg
        bg.line.fill.background()
        add_corner_logo(slide)

        tbox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(10.0), Inches(0.8))
        tbox.text_frame.paragraphs[0].text = title
        tbox.text_frame.paragraphs[0].font.size = Pt(24)
        tbox.text_frame.paragraphs[0].font.bold = True
        tbox.text_frame.paragraphs[0].font.color.rgb = c_wine
        tbox.text_frame.paragraphs[0].font.name = font_bold

        slide.shapes.add_picture(img_func, Inches(0.8), Inches(1.3), width=Inches(11.7))

    # СЛАЙД С ТАБЛИЦЕЙ ДЕТАЛИЗАЦИИ
    slide_table = prs.slides.add_slide(blank_layout)
    bg_t = slide_table.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg_t.fill.solid()
    bg_t.fill.fore_color.rgb = c_light_bg
    bg_t.line.fill.background()
    add_corner_logo(slide_table)

    title_box_t = slide_table.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(10.0), Inches(0.8))
    title_box_t.text_frame.paragraphs[0].text = "Детализация денежных потоков по месяцам"
    title_box_t.text_frame.paragraphs[0].font.size = Pt(26)
    title_box_t.text_frame.paragraphs[0].font.bold = True
    title_box_t.text_frame.paragraphs[0].font.color.rgb = c_wine
    title_box_t.text_frame.paragraphs[0].font.name = font_bold

    rows = min(period + 1, 13)
    cols = 5
    table_shape = slide_table.shapes.add_table(rows, cols, Inches(0.8), Inches(1.3), Inches(11.7), Inches(5.4))
    table = table_shape.table
    
    table.columns[0].width = Inches(2.2)
    table.columns[1].width = Inches(2.4)
    table.columns[2].width = Inches(2.4)
    table.columns[3].width = Inches(2.4)
    table.columns[4].width = Inches(2.3)

    headers = ["Месяц", "Выручка (руб)", "Поступления (руб)", "Выплаты (руб)", "Остаток ДС (руб)"]
    for col_idx, text in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = c_wine
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = c_white
            p.font.name = font_bold
            p.alignment = PP_ALIGN.CENTER

    for i in range(rows - 1):
        row_data = [
            str(x_labels[i]),
            f"{rev[i]:,.0f}".replace(",", " "),
            f"{inflows[i]:,.0f}".replace(",", " "),
            f"{outflows[i]:,.0f}".replace(",", " "),
            f"{cash_balance[i]:,.0f}".replace(",", " ")
        ]
        for col_idx, val in enumerate(row_data):
            cell = table.cell(i+1, col_idx)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = c_white if i % 2 == 0 else RGBColor(240, 235, 230)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = c_dark
                p.font.name = font_regular
                p.alignment = PP_ALIGN.CENTER if col_idx == 0 else PP_ALIGN.RIGHT

    ppt_io = io.BytesIO()
    prs.save(ppt_io)
    ppt_io.seek(0)
    return ppt_io

def send_report_to_email(to_email, pptx_bytes):
    smtp_server = os.getenv("SMTP_SERVER", "smtp.yandex.ru")
    smtp_port = int(os.getenv("SMTP_PORT", 465))
    sender_email = os.getenv("SMTP_USER", "finance@krayvin.ru")
    sender_password = os.getenv("SMTP_PASSWORD", "ваш_пароль_приложения")

    msg = MIMEMultipart()
    msg['From'] = sender_email
    msg['To'] = to_email
    msg['Subject'] = "🍷 КРАЙВИН: Полный финансовый отчёт и модель денежных потоков"

    body = "Здравствуйте!\n\nВо вложении находится обновленный инвестиционный отчет компании КРАЙВИН в формате PowerPoint со всеми новыми аналитическими графиками.\n\nС уважением,\nФинансовый департамент КРАЙВИН"
    msg.attach(MIMEText(body, 'plain'))

    attachment = MIMEBase('application', 'vnd.openxmlformats-officedocument.presentationml.presentation')
    attachment.set_payload(pptx_bytes.getvalue())
    encoders.encode_base64(attachment)
    attachment.add_header('Content-Disposition', 'attachment', filename="Kraivin_Investment_Report.pptx")
    msg.attach(attachment)

    try:
        server = smtplib.SMTP_SSL(smtp_server, smtp_port)
        server.login(sender_email, sender_password)
        server.sendmail(sender_email, to_email, msg.as_string())
        server.quit()
        return True
    except Exception as e:
        print(f"SMTP Error: {e}")
        return False
