import streamlit as st
import pandas as pd
import plotly.graph_objects as go
import numpy as np
import os
import io

# Импорт библиотек для генерации PowerPoint
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Настройка страницы
st.set_page_config(
    page_title="КРАЙВИН - анализ экономической эффективности", 
    page_icon="🍷", 
    layout="wide"
)

st.title("КРАЙВИН: Анализ денежных потоков и рентабельности")
st.markdown("Интерактивная финансовая модель для сценарного анализа кассовых разрывов.")

# --- БОКОВАЯ ПАНЕЛЬ (ВВОД ДАННЫХ) ---

logo_path = "КРАЙВИН лого винный квадрат.png"
if os.path.exists(logo_path):
    st.sidebar.image(logo_path, use_container_width=True)
else:
    st.sidebar.info("💡 Загрузите файл 'КРАЙВИН лого винный квадрат.png' в папку с кодом для отображения логотипа.")

st.sidebar.header("Параметры модели")

ru_months_full = ['Январь', 'Февраль', 'Март', 'Апрель', 'Май', 'Июнь', 'Июль', 'Август', 'Сентябрь', 'Октябрь', 'Ноябрь', 'Декабрь']
col_m, col_y = st.sidebar.columns(2)
start_month_idx = col_m.selectbox("Месяц старта", range(12), format_func=lambda x: ru_months_full[x])
start_year = col_y.selectbox("Год старта", [2026, 2027])

margin_pct = st.sidebar.slider("Маржинальность (%)", min_value=10, max_value=50, value=20, step=1)
period = st.sidebar.selectbox("Горизонт планирования (мес)", [6, 12, 18, 24])

st.sidebar.subheader("Стартовый капитал и закупки")
initial_purchase = st.sidebar.number_input("Первоначальная закупка товара (руб)", value=5_000_000, step=500_000)
initial_cash_buffer = st.sidebar.number_input("Стартовый денежный буфер (на счете)", value=2_000_000, step=500_000)

st.sidebar.subheader("Динамика продаж")
aov = st.sidebar.number_input("Средняя сумма заказа (руб)", value=150_000, step=10_000)
start_orders = st.sidebar.number_input("Заказов в 1-й месяц (шт)", value=40, step=1)
orders_growth = st.sidebar.slider("Ежемесячный прирост заказов (%)", 0, 100, 15, step=1)
scale_factor = st.sidebar.slider("Коэффициент масштабирования продаж", 0.5, 3.0, 1.0, 0.1)

st.sidebar.subheader("Команда и расходы")
monthly_fot = st.sidebar.number_input("ФОТ в месяц (руб)", value=500_000, step=50_000)

st.sidebar.subheader("Работа с поставщиками")
prepayment_pct = st.sidebar.slider("Предоплата поставщикам (%)", 0, 100, 50, step=10)
delay_days = st.sidebar.slider("Отсрочка на остаток (дней)", 0, 90, 40, step=5)

st.sidebar.subheader("Факторинг")
factoring_share = st.sidebar.slider("Доля выручки в факторинге (%)", 0, 100, 50, step=10)
factoring_advance = st.sidebar.slider("Аванс от фактора (%)", 50, 100, 80, step=5)

st.sidebar.subheader("Условия с покупателями")
customer_delay_days = st.sidebar.slider("Отсрочка платежа покупателям (дней)", 0, 120, 70, step=5)

# --- РАСЧЕТНАЯ ЧАСТЬ (МАТЕМАТИКА) ---

ru_months_short = ['Янв', 'Фев', 'Мар', 'Апр', 'Май', 'Июн', 'Июл', 'Авг', 'Сен', 'Окт', 'Ноя', 'Дек']
x_labels = []
for i in range(period):
    m_idx = (start_month_idx + i) % 12
    y_offset = (start_month_idx + i) // 12
    x_labels.append(f"{ru_months_short[m_idx]} {start_year + y_offset}")

orders = np.zeros(period)
rev = np.zeros(period)

for i in range(period):
    if i == 0:
        orders[i] = start_orders * scale_factor
    else:
        orders[i] = orders[i-1] * (1 + (orders_growth / 100))
    rev[i] = orders[i] * aov

cogs_pct = 1 - (margin_pct / 100)
cogs_no_vat = rev * cogs_pct
cogs_vat = cogs_no_vat * 1.2

delay_months_suppliers = max(1, int(round(delay_days / 30))) if delay_days > 0 else 0
cogs_payments = np.zeros(period)

for i in range(period):
    cogs_payments[i] += cogs_vat[i] * (prepayment_pct / 100)
    if i + delay_months_suppliers < period:
        cogs_payments[i + delay_months_suppliers] += cogs_vat[i] * ((100 - prepayment_pct) / 100)

initial_prep = initial_purchase * (prepayment_pct / 100)
initial_post = initial_purchase * ((100 - prepayment_pct) / 100)

cogs_payments[0] += initial_prep
if delay_months_suppliers < period:
    cogs_payments[delay_months_suppliers] += initial_post

customer_delay_months = max(0, int(round(customer_delay_days / 30)))

inflows = np.zeros(period)
for i in range(period):
    inflows[i] += rev[i] * 1.2 * (factoring_share / 100) * (factoring_advance / 100)
    
    target_month = i + customer_delay_months
    if target_month < period:
        inflows[target_month] += rev[i] * 1.2 * ((100 - factoring_share) / 100)
        inflows[target_month] += rev[i] * 1.2 * (factoring_share / 100) * ((100 - factoring_advance) / 100)

base_other_opex = 150_000
opex = np.full(period, base_other_opex + monthly_fot)
for i in range(6, period):
    opex[i] = (base_other_opex * 1.2) + monthly_fot

taxes_and_commissions = rev * 0.05

outflows = cogs_payments + opex + taxes_and_commissions
net_cf = inflows - outflows

cum_cf = np.cumsum(net_cf)
cash_balance = cum_cf + initial_cash_buffer

# --- KPI МЕТРИКИ ---
max_deficit = min(min(cum_cf), 0)
net_profit = sum(rev * (margin_pct / 100)) - sum(opex) - sum(taxes_and_commissions) - (initial_purchase * 0.15)
roi = (net_profit / sum(rev)) * 100 if sum(rev) > 0 else 0

def format_rub(val):
    return f"{val:,.0f}".replace(",", " ") + " руб."

col1, col2, col3, col4 = st.columns(4)
col1.metric(f"Выручка (за {period} мес)", format_rub(sum(rev)))
col2.metric("Макс. кассовый разрыв", format_rub(max_deficit))
col3.metric("Чистая прибыль", format_rub(net_profit))
col4.metric("Рентабельность по ЧП", f"{roi:.1f}%")

st.divider()

# --- ПРОФЕССИОНАЛЬНЫЙ ГЕНЕРАТОР ПРЕЗЕНТАЦИЙ (PPTX) ---
def generate_professional_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333) # Широкий формат 16:9
    prs.slide_height = Inches(7.5)
    
    # Фирменная палитра
    c_wine = RGBColor(100, 42, 56)      # #642A38 (Винный)
    c_sand = RGBColor(227, 194, 147)   # #E3C293 (Песочный)
    c_dark = RGBColor(30, 30, 30)      # Темно-серый текст
    c_light_bg = RGBColor(248, 246, 244) # Светлый фон слайдов
    c_white = RGBColor(255, 255, 255)
    c_card_bg = RGBColor(255, 255, 255)
    c_card_border = RGBColor(220, 210, 205)

    blank_layout = prs.slide_layouts[6] # Пустой слайд для кастомной верстки

    # --- СЛАЙД 1: ТИТУЛЬНЫЙ (ПРЕМИУМ СТИЛЬ) ---
    slide1 = prs.slides.add_slide(blank_layout)
    
    # Заливка фона титульного слайда винным цветом
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = c_wine
    bg1.line.fill.background()

    # Декоративный песочный акцент (плашка слева)
    accent1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.0), Inches(0.15), Inches(3.5))
    accent1.fill.solid()
    accent1.fill.fore_color.rgb = c_sand
    accent1.line.fill.background()

    # Текст заголовка
    txBox = slide1.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(11.0), Inches(3.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "КРАЙВИН"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = c_white
    p.font.name = "Arial"

    p2 = tf.add_paragraph()
    p2.text = "Финансовая модель и анализ денежных потоков"
    p2.font.size = Pt(28)
    p2.font.color.rgb = c_sand
    p2.font.name = "Arial"
    p2.space_before = Pt(10)

    p3 = tf.add_paragraph()
    p3.text = f"Горизонт планирования: {period} мес.  |  Старт: {ru_months_full[start_month_idx]} {start_year}"
    p3.font.size = Pt(16)
    p3.font.color.rgb = c_white
    p3.font.name = "Arial"
    p3.space_before = Pt(40)


    # --- СЛАЙД 2: КЛЮЧЕВЫЕ ПОКАЗАТЕЛИ (КАРТОЧКИ KPI) ---
    slide2 = prs.slides.add_slide(blank_layout)
    
    # Фоновая подложка
    bg2 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg2.fill.solid()
    bg2.fill.fore_color.rgb = c_light_bg
    bg2.line.fill.background()

    # Шапка слайда
    title_box = slide2.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.0), Inches(1.0))
    tf2 = title_box.text_frame
    p_t = tf2.paragraphs[0]
    p_t.text = "Ключевые финансовые результаты"
    p_t.font.size = Pt(32)
    p_t.font.bold = True
    p_t.font.color.rgb = c_wine
    p_t.font.name = "Arial"

    # Массив KPI для карточек
    kpis = [
        ("Суммарная выручка", format_rub(sum(rev))),
        ("Макс. кассовый разрыв", format_rub(max_deficit)),
        ("Чистая прибыль", format_rub(net_profit)),
        ("Рентабельность по ЧП", f"{roi:.1f}%"),
        ("Начальный буфер ДС", format_rub(initial_cash_buffer)),
        ("Первоначальная закупка", format_rub(initial_purchase))
    ]

    # Рисуем сетку карточек 3х2
    card_w, card_h = Inches(3.6), Inches(2.2)
    start_x, start_y = Inches(0.8), Inches(1.8)
    gap_x, gap_y = Inches(0.4), Inches(0.4)

    for idx, (label, val) in enumerate(kpis):
        col_idx = idx % 3
        row_idx = idx // 3
        x = start_x + col_idx * (card_w + gap_x)
        y = start_y + row_idx * (card_h + gap_y)

        # Подложка карточки
        card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = c_white
        card.line.color.rgb = c_card_border
        card.line.width = Pt(1)

        # Текст внутри карточки
        tf_card = card.text_frame
        tf_card.word_wrap = True
        tf_card.margin_left = Inches(0.2)
        tf_card.margin_top = Inches(0.2)
        
        p_lbl = tf_card.paragraphs[0]
        p_lbl.text = label.upper()
        p_lbl.font.size = Pt(12)
        p_lbl.font.color.rgb = RGBColor(120, 120, 120)
        p_lbl.font.bold = True

        p_val = tf_card.add_paragraph()
        p_val.text = val
        p_val.font.size = Pt(22)
        p_val.font.color.rgb = c_wine
        p_val.font.bold = True
        p_val.space_before = Pt(10)


    # --- СЛАЙД 3: ТАБЛИЦА ДЕТАЛИЗАЦИИ ---
    slide3 = prs.slides.add_slide(blank_layout)
    
    bg3 = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg3.fill.solid()
    bg3.fill.fore_color.rgb = c_light_bg
    bg3.line.fill.background()

    title_box3 = slide3.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.0), Inches(0.8))
    tf3 = title_box3.text_frame
    p_t3 = tf3.paragraphs[0]
    p_t3.text = "Детализация денежных потоков по месяцам"
    p_t3.font.size = Pt(28)
    p_t3.font.bold = True
    p_t3.font.color.rgb = c_wine

    # Создание таблицы
    rows = min(period + 1, 13) # Ограничим до 12 месяцев на слайд для читаемости
    cols = 5
    t_left, t_top, t_width, t_height = Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.2)
    
    table_shape = slide3.shapes.add_table(rows, cols, t_left, t_top, t_width, t_height)
    table = table_shape.table
    
    # Настройка ширины колонок
    table.columns[0].width = Inches(2.2)
    table.columns[1].width = Inches(2.4)
    table.columns[2].width = Inches(2.4)
    table.columns[3].width = Inches(2.4)
    table.columns[4].width = Inches(2.3)

    headers = ["Месяц", "Выручка (руб)", "Поступления (руб)", "Выплаты (руб)", "Остаток ДС (руб)"]
    for col_idx, text in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = c_wine
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = c_white
            p.alignment = PP_ALIGN.CENTER

    for i in range(rows - 1):
        row_data = [
            str(x_labels[i]),
            f"{rev[i]:,.0f}".replace(",", " "),
            f"{inflows[i]:,.0f}".replace(",", " "),
            f"{outflows[i]:,.0f}".replace(",", " "),
            f"{cash_balance[i]:,.0f}".replace(",", " ")
        ]
        for col_idx, val in enumerate(row_data):
            cell = table.cell(i+1, col_idx)
            cell.text = val
            cell.fill.solid()
            if i % 2 == 0:
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
            else:
                cell.fill.fore_color.rgb = RGBColor(240, 235, 230)
                
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.color.rgb = c_dark
                p.alignment = PP_ALIGN.CENTER if col_idx == 0 else PP_ALIGN.RIGHT

    # Сохранение в буфер
    ppt_io = io.BytesIO()
    prs.save(ppt_io)
    ppt_io.seek(0)
    return ppt_io

# Кнопка скачивания презентации в интерфейсе
st.sidebar.divider()
st.sidebar.subheader("Экспорт отчета")
if st.sidebar.button("📥 Скачать профессиональную презентацию (PPTX)"):
    pptx_data = generate_professional_pptx()
    st.sidebar.download_button(
        label="💾 Нажмите для сохранения файла",
        data=pptx_data,
        file_name="Kraivin_Professional_Report.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )

# --- ВИЗУАЛИЗАЦИЯ (ГРАФИКИ PLOTLY) ---
st.subheader("Динамика ликвидности и остаток средств")

fig1 = go.Figure()
fig1.add_trace(go.Scatter(
    x=x_labels, 
    y=cash_balance, 
    mode='lines+markers', 
    name='Остаток ДС',
    line=dict(color='#642A38', width=3),
    fill='tozeroy',
    fillcolor='rgba(100, 42, 56, 0.1)',
    hovertemplate='%{y:,.0f} руб.<extra></extra>'
))
fig1.add_hline(y=0, line_dash="dash", line_color="red", annotation_text="Дефицит")
fig1.update_layout(
    xaxis_title="Месяц", 
    yaxis_title="Рубли", 
    hovermode="x unified",
    separators=", "
)
fig1.update_yaxes(tickformat=",.0f")
st.plotly_chart(fig1, use_container_width=True)


st.subheader("Структура месячного денежного потока")
fig2 = go.Figure()
fig2.add_trace(go.Bar(
    x=x_labels, y=inflows, name='Поступления', marker_color='#E3C293',
    hovertemplate='%{y:,.0f} руб.<extra></extra>'
))
fig2.add_trace(go.Bar(
    x=x_labels, y=-outflows, name='Выплаты', marker_color='#642A38',
    hovertemplate='%{y:,.0f} руб.<extra></extra>'
))
fig2.add_trace(go.Scatter(
    x=x_labels, y=net_cf, name='Чистый поток', marker_color='#B88645',
    mode='lines+markers',
    hovertemplate='%{y:,.0f} руб.<extra></extra>'
))

fig2.update_layout(
    barmode='relative', 
    xaxis_title="Месяц", 
    yaxis_title="Рубли", 
    hovermode="x unified",
    separators=", "
)
fig2.update_yaxes(tickformat=",.0f")
st.plotly_chart(fig2, use_container_width=True)
